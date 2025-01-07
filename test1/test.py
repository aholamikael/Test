import os
import json
from azure.storage.blob import BlobServiceClient
from azure.identity import DefaultAzureCredential
from pptx import Presentation

def create_blob_service_client(storage_account_url: str, credential):

    return BlobServiceClient(account_url=storage_account_url, credential=credential)

def download_blob_to_local(blob_client, blob_name: str, local_path: str):

    blob_data = blob_client.download_blob(blob_name)
    with open(local_path, "wb") as f:
        f.write(blob_data.readall())

def extract_text_from_pptx(local_pptx_path: str):

    ppt = Presentation(local_pptx_path)
    slides_text = []
    for slide in ppt.slides:
        slide_content = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_content.append(shape.text)
        slides_text.append(slide_content)
    
    return json.dumps(slides_text, indent=1)

def upload_json_to_blob(blob_client, json_data: str, json_blob_name: str):

    blob_client.upload_blob(json_data, overwrite=True)

def process_pptx_files():

    storage_account_url = "https://fall24skillionidatastore.blob.core.windows.net/"
    original_container_name = "data-visualisation"
    conversion_container_name = "testconversions"

    credential = DefaultAzureCredential()
    blob_service_client = create_blob_service_client(storage_account_url, credential)
    original_container_client = blob_service_client.get_container_client(original_container_name)
    conversion_container_client = blob_service_client.get_container_client(conversion_container_name)

    try:
        pptx_blobs = original_container_client.list_blobs()

        for pptx_blob in pptx_blobs:
            if pptx_blob.name.endswith(".pptx"):
                pptx_blob_name = pptx_blob.name
                print(f"Processing file: {pptx_blob_name}")

                local_pptx_path = "/tmp/" + pptx_blob_name

                download_blob_to_local(original_container_client, pptx_blob_name, local_pptx_path)

                json_data = extract_text_from_pptx(local_pptx_path)

                json_blob_name = pptx_blob_name.replace(".pptx", ".json")

                upload_json_to_blob(conversion_container_client, json_data, json_blob_name)

                os.remove(local_pptx_path)

        print("Successfully processed all PPTX files.")

    except Exception as e:
        print(f"Error processing files: {str(e)}")
